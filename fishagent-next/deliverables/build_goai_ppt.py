from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET = ROOT / "deliverables" / "assets"
OUT = ROOT / "deliverables" / "智渔Agent-2.0_GOAI无界应用评委路演.pptx"
FONT = "Noto Sans CJK SC"
W, H = 13.333, 7.5


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NAVY = rgb("0B2A56")
BLUE = rgb("0B66C3")
CYAN = rgb("26C6E6")
MINT = rgb("41D19A")
YELLOW = rgb("FFD13D")
CORAL = rgb("FF6B5E")
PURPLE = rgb("7255E8")
INK = rgb("14233B")
MUTED = rgb("5D6B80")
LINE = rgb("D7E1ED")
PALE = rgb("F4F7FB")
PALE_BLUE = rgb("EAF3FF")
PALE_MINT = rgb("EAF9F3")
PALE_YELLOW = rgb("FFF8DF")
PALE_CORAL = rgb("FFF0ED")
WHITE = rgb("FFFFFF")


def x(value: float) -> int:
    return Inches(value)


def add_shape(slide, shape_type, left, top, width, height, fill=WHITE, line=LINE, radius=True):
    shape = slide.shapes.add_shape(shape_type, x(left), x(top), x(width), x(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    if radius and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        shape.adjustments[0] = 0.12
    return shape


def card(slide, left, top, width, height, fill=WHITE, line=LINE):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line, True)


def text(slide, value, left, top, width, height, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT,
         valign=MSO_ANCHOR.TOP, font=FONT, margin=0.08, italic=False):
    box = slide.shapes.add_textbox(x(left), x(top), x(width), x(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = x(margin)
    frame.margin_right = x(margin)
    frame.margin_top = x(margin)
    frame.margin_bottom = x(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def rich_text(slide, runs, left, top, width, height, size=16, color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x(left), x(top), x(width), x(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = x(0.08)
    frame.margin_right = x(0.08)
    frame.margin_top = x(0.08)
    frame.margin_bottom = x(0.08)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for value, run_size, run_color, bold in runs:
        run = paragraph.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(run_size or size)
        run.font.bold = bold
        run.font.color.rgb = run_color or color
    return box


def bullet_list(slide, items, left, top, width, height, size=14, color=INK, gap=4, bullet_color=None):
    box = slide.shapes.add_textbox(x(left), x(top), x(width), x(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = x(0.08)
    frame.margin_right = x(0.08)
    frame.margin_top = x(0.05)
    frame.margin_bottom = x(0.05)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"·  {item}"
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(gap)
    return box


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2, arrow=False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(x1), x(y1), x(x2), x(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    if arrow:
        connector.line.end_arrowhead = True
    return connector


def add_picture(slide, path, left, top, width, height):
    return slide.shapes.add_picture(str(path), x(left), x(top), width=x(width), height=x(height))


def crop(path: Path, target: Path, box: tuple[int, int, int, int]) -> Path:
    if not target.exists() or target.stat().st_mtime < path.stat().st_mtime:
        image = Image.open(path)
        image.crop(box).save(target)
    return target


def source_note(slide, value, left=0.6, top=7.17, width=11.7, color=MUTED):
    text(slide, value, left, top, width, 0.2, size=7.5, color=color, margin=0)


def footer(slide, number, dark=False):
    color = rgb("B8CAE0") if dark else MUTED
    text(slide, f"智渔Agent 2.0  /  GOAI Boundless Agents                                      {number:02d}",
         0.6, 7.18, 12.1, 0.18, size=7.5, color=color, margin=0)


def base(slide, title, kicker=None, dark=False, number=None):
    if dark:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = NAVY
        text(slide, "GOAI  /  BOUNDLESS AGENTS", 0.62, 0.34, 4.4, 0.22, size=8, color=CYAN, bold=True, margin=0)
        if kicker:
            text(slide, kicker.upper(), 0.62, 0.72, 5.6, 0.22, size=8, color=YELLOW, bold=True, margin=0)
        text(slide, title, 0.62, 0.96, 11.8, 0.62, size=27, color=WHITE, bold=True, margin=0)
    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PALE
        text(slide, "GOAI  /  BOUNDLESS AGENTS", 0.62, 0.34, 4.4, 0.22, size=8, color=BLUE, bold=True, margin=0)
        if kicker:
            text(slide, kicker.upper(), 0.62, 0.72, 5.6, 0.22, size=8, color=PURPLE, bold=True, margin=0)
        text(slide, title, 0.62, 0.96, 11.8, 0.62, size=27, color=INK, bold=True, margin=0)
    if number is not None:
        footer(slide, number, dark=dark)


def dot_field(slide, left, top, width, height, dark=True):
    colors = [CYAN, MINT, YELLOW, CORAL, PURPLE, BLUE]
    positions = [
        (0.1, 0.1, 0.07), (0.36, 0.18, 0.15), (0.68, 0.03, 0.08), (0.82, 0.16, 0.2),
        (0.16, 0.42, 0.2), (0.48, 0.4, 0.11), (0.76, 0.38, 0.14), (0.88, 0.54, 0.08),
        (0.28, 0.68, 0.09), (0.55, 0.72, 0.22), (0.78, 0.74, 0.1), (0.12, 0.88, 0.16),
    ]
    for index, (px, py, size) in enumerate(positions):
        add_shape(slide, MSO_SHAPE.OVAL, left + px * width, top + py * height,
                  size, size, colors[index % len(colors)], colors[index % len(colors)], False)


def metric_card(slide, left, top, width, value, label, accent, sub=None):
    card(slide, left, top, width, 1.05, WHITE, LINE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, 0.06, 1.05, accent, accent, False)
    text(slide, value, left + 0.18, top + 0.16, width - 0.3, 0.36, size=24, color=accent, bold=True, margin=0)
    text(slide, label, left + 0.18, top + 0.57, width - 0.3, 0.21, size=10.5, color=INK, bold=True, margin=0)
    if sub:
        text(slide, sub, left + 0.18, top + 0.79, width - 0.3, 0.16, size=7.5, color=MUTED, margin=0)


def pill(slide, label, left, top, width, fill, color=INK):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, 0.32, fill, fill, True)
    text(slide, label, left, top + 0.035, width, 0.21, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)


def numbered_callout(slide, number, title, body, left, top, width, fill=WHITE, accent=BLUE):
    card(slide, left, top, width, 0.72, fill, LINE)
    add_shape(slide, MSO_SHAPE.OVAL, left + 0.12, top + 0.14, 0.4, 0.4, accent, accent, False)
    text(slide, str(number), left + 0.12, top + 0.19, 0.4, 0.22, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    text(slide, title, left + 0.62, top + 0.1, width - 0.74, 0.2, size=11, color=INK, bold=True, margin=0)
    text(slide, body, left + 0.62, top + 0.34, width - 0.74, 0.26, size=8.5, color=MUTED, margin=0)


def make_slides():
    monitor_top = crop(ASSET / "fishagent-monitor.png", ASSET / "fishagent-monitor-top.png", (0, 0, 1440, 960))
    report_top = crop(ASSET / "fishagent-reports.png", ASSET / "fishagent-reports-top.png", (0, 0, 1440, 900))

    presentation = Presentation()
    presentation.slide_width = x(W)
    presentation.slide_height = x(H)
    blank = presentation.slide_layouts[6]

    # 1. Cover
    slide = presentation.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    dot_field(slide, 8.0, 0.8, 4.7, 5.5)
    add_picture(slide, ASSET / "goai-logo.png", 10.65, 0.35, 1.45, 0.49)
    text(slide, "GOAI 无界应用", 0.72, 0.62, 4.8, 0.3, size=12, color=YELLOW, bold=True, margin=0)
    text(slide, "智渔Agent 2.0", 0.68, 1.35, 7.4, 0.84, size=38, color=WHITE, bold=True, margin=0)
    text(slide, "面向水产养殖的可验证闭环 Agent", 0.72, 2.3, 6.9, 0.48, size=22, color=CYAN, bold=True, margin=0)
    text(slide, "把感知变成决策，把决策变成动作，把动作变成可验证结果。", 0.74, 3.05, 6.6, 0.36, size=14, color=rgb("D4E5F7"), margin=0)
    pill(slide, "Boundless Agents 赛道", 0.74, 4.05, 2.15, BLUE, WHITE)
    pill(slide, "真实行业场景", 3.02, 4.05, 1.55, MINT, NAVY)
    pill(slide, "可运行 Demo", 4.72, 4.05, 1.35, YELLOW, NAVY)
    card(slide, 0.72, 5.2, 5.6, 0.95, rgb("143A6B"), rgb("2C6296"))
    text(slide, "参赛团队：__________   成员：__________", 0.95, 5.42, 5.15, 0.23, size=12, color=WHITE, margin=0)
    text(slide, "决赛路演版 · 2026", 0.95, 5.78, 4.6, 0.19, size=9, color=rgb("B8CAE0"), margin=0)
    text(slide, "评委先看闭环，再看技术。", 8.0, 6.45, 4.45, 0.28, size=14, color=WHITE, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    text(slide, "可运行 · 可演示 · 可复制", 8.0, 6.82, 4.45, 0.22, size=10, color=YELLOW, align=PP_ALIGN.RIGHT, margin=0)

    # 2. Evaluation fit
    slide = presentation.slides.add_slide(blank)
    base(slide, "先给评委一个可验证的答案", "01 / 项目定位", number=2)
    text(slide, "智渔Agent不是把聊天框接到传感器上，而是把养殖异常变成一条可追溯、可复核、可交接的任务链。", 0.68, 1.8, 12, 0.42, size=15, color=INK, bold=True, margin=0)
    metric_card(slide, 0.68, 2.45, 2.25, "4", "养殖池塘", BLUE, "B-01 至 B-04")
    metric_card(slide, 3.07, 2.45, 2.25, "28", "模拟设备总数", MINT, "正常 28/28 · 故障 Demo 27/28")
    metric_card(slide, 5.46, 2.45, 2.25, "7", "核心传感器指标", PURPLE, "水质全量巡检")
    metric_card(slide, 7.85, 2.45, 2.25, "MQTT", "设备消息总线", CORAL, "上报与下发统一走消息")
    card(slide, 10.25, 2.45, 2.4, 1.05, PALE_YELLOW, YELLOW)
    text(slide, "1 条", 10.45, 2.62, 1.95, 0.35, size=24, color=CORAL, bold=True, margin=0)
    text(slide, "完整任务闭环", 10.45, 3.02, 1.95, 0.2, size=10.5, color=INK, bold=True, margin=0)
    text(slide, "告警 → 决策 → 执行 → 复核", 10.45, 3.25, 1.95, 0.16, size=7.3, color=MUTED, margin=0)
    card(slide, 0.68, 3.9, 5.25, 2.48, WHITE, LINE)
    text(slide, "与 GOAI 赛道定位的交集", 0.95, 4.16, 4.7, 0.3, size=16, color=INK, bold=True, margin=0)
    bullet_list(slide, [
        "真实行业：水产养殖异常处置，而非泛问答",
        "多轮交互：聊天、主动巡塘、告警复核相互衔接",
        "工具调用：知识库、天气、MQTT、报告交付",
        "多模态：传感器 + 水面 / 水下图像 + 天气上下文",
    ], 0.96, 4.62, 4.6, 1.45, size=11.5)
    card(slide, 6.2, 3.9, 6.45, 2.48, WHITE, LINE)
    text(slide, "按手册 10.1 逐项准备证据", 6.48, 4.16, 5.7, 0.3, size=16, color=INK, bold=True, margin=0)
    rubric = [
        ("行业场景价值", "25%", "低溶氧、设备离线、人工接管"),
        ("Agent 能力与任务闭环", "25%", "结构化决策 + Skill + 复核"),
        ("产品体验与 Demo", "20%", "页面、流式轨迹、告警流程"),
        ("技术实现深度", "15%", "CrewAI / MQTT / 持久化 / 测试"),
        ("安全合规与可追溯", "10%", "策略门、人工任务、审计"),
        ("开放 / 复用贡献", "5%", "代码、文档、模拟数据规格"),
    ]
    for idx, (label, weight, evidence) in enumerate(rubric):
        y = 4.62 + idx * 0.27
        text(slide, label, 6.48, y, 1.72, 0.19, size=8.8, color=INK, bold=True, margin=0)
        text(slide, weight, 8.15, y, 0.48, 0.19, size=8.8, color=BLUE, bold=True, margin=0)
        text(slide, evidence, 8.72, y, 3.55, 0.19, size=8.8, color=MUTED, margin=0)
    source_note(slide, "依据：GOAI《无界应用｜Boundless Agents 参赛手册》10.1 通用评审维度；官网：goaihz.com")

    # 3. Pain and users
    slide = presentation.slides.add_slide(blank)
    base(slide, "真实场景：异常处置不是一句建议", "02 / 行业价值", number=3)
    text(slide, "典型场景：夜间或无人值守时，B-01 溶解氧跌破安全线；养殖户要在风险扩大前完成判断、动作和复核，而不是只收到一条告警。", 0.68, 1.78, 12, 0.45, size=14, color=INK, bold=True, margin=0)
    users = [
        ("养殖户", "关心损失与收益", "少死鱼、少误操作、少熬夜", PALE_BLUE, BLUE),
        ("现场操作员", "关心动作与交接", "任务清楚、失败可接管、结果可复核", PALE_MINT, MINT),
        ("管理者", "关心证据与复制", "趋势、审计、报告、跨池塘复用", PALE_YELLOW, YELLOW),
    ]
    for idx, (title, sub, body, fill, accent) in enumerate(users):
        left = 0.68 + idx * 4.17
        card(slide, left, 2.5, 3.8, 1.32, fill, accent)
        text(slide, title, left + 0.22, 2.72, 1.35, 0.24, size=15, color=INK, bold=True, margin=0)
        text(slide, sub, left + 1.53, 2.74, 1.98, 0.2, size=9.5, color=accent, bold=True, margin=0)
        text(slide, body, left + 0.22, 3.13, 3.35, 0.24, size=10.5, color=MUTED, margin=0)
    text(slide, "传统流程", 0.8, 4.35, 1.2, 0.24, size=12, color=MUTED, bold=True, margin=0)
    old_steps = ["传感器报警", "人工判断", "打电话找人", "打开设备", "等待结果"]
    for idx, item in enumerate(old_steps):
        left = 2.0 + idx * 2.05
        card(slide, left, 4.2, 1.6, 0.54, WHITE, LINE)
        text(slide, item, left, 4.36, 1.6, 0.18, size=9.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER, margin=0)
        if idx < len(old_steps) - 1:
            line(slide, left + 1.62, 4.47, left + 1.96, 4.47, LINE, 1.5, True)
    text(slide, "智渔Agent", 0.8, 5.25, 1.2, 0.24, size=12, color=BLUE, bold=True, margin=0)
    new_steps = ["主动取数", "Agent 研判", "Skill 校验", "MQTT 执行", "延时复核"]
    for idx, item in enumerate(new_steps):
        left = 2.0 + idx * 2.05
        fill = PALE_BLUE if idx < 2 else PALE_MINT
        accent = BLUE if idx < 2 else MINT
        card(slide, left, 5.1, 1.6, 0.54, fill, accent)
        text(slide, item, left, 5.26, 1.6, 0.18, size=9.5, color=INK, bold=True, align=PP_ALIGN.CENTER, margin=0)
        if idx < len(new_steps) - 1:
            line(slide, left + 1.62, 5.37, left + 1.96, 5.37, accent, 1.5, True)
    card(slide, 0.68, 6.12, 12.0, 0.42, NAVY, NAVY)
    text(slide, "关键差异：模型负责理解与决策，Skill 负责可行性，MQTT 负责设备边界，复核 Agent 负责判断处置是否真的生效。", 0.9, 6.22, 11.55, 0.2, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    # 4. Closed loop
    slide = presentation.slides.add_slide(blank)
    base(slide, "一个完整闭环：以 B-01 低溶氧为例", "03 / Agent 闭环", number=4)
    text(slide, "示例输入：B-01 溶解氧 2.8 mg/L，安全下限 4.0 mg/L。系统不止报错，而是把处理过程跑完。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    steps = [
        ("01", "任务输入", "MQTT 传感器\nDO 2.8 mg/L", BLUE, PALE_BLUE),
        ("02", "意图理解", "sensor-monitor\n确认新鲜读数", CYAN, PALE_BLUE),
        ("03", "任务规划", "action-planning\n形成 action", PURPLE, rgb("F0EDFF")),
        ("04", "Skill 校验", "能力 / 风险 /\n幂等 / 健康", YELLOW, PALE_YELLOW),
        ("05", "设备执行", "execution-agent\n调用设备 Skill", MINT, PALE_MINT),
        ("06", "延时复核", "verification-agent\n等待复核", CORAL, PALE_CORAL),
        ("07", "结果交付", "关闭告警 /\n人工接管 / 报告", BLUE, PALE_BLUE),
    ]
    for idx, (number, title, body, accent, fill) in enumerate(steps):
        left = 0.68 + idx * 1.78
        card(slide, left, 2.55, 1.5, 1.56, fill, accent)
        add_shape(slide, MSO_SHAPE.OVAL, left + 0.54, 2.72, 0.42, 0.42, accent, accent, False)
        text(slide, number, left + 0.54, 2.82, 0.42, 0.16, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        text(slide, title, left + 0.1, 3.27, 1.3, 0.2, size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER, margin=0)
        text(slide, body, left + 0.1, 3.58, 1.3, 0.36, size=8.5, color=MUTED, align=PP_ALIGN.CENTER, margin=0)
        if idx < len(steps) - 1:
            line(slide, left + 1.52, 3.33, left + 1.72, 3.33, accent, 1.4, True)
    card(slide, 0.68, 4.52, 5.6, 1.42, WHITE, LINE)
    text(slide, "自动处置成功", 0.95, 4.78, 2.0, 0.26, size=15, color=MINT, bold=True, margin=0)
    text(slide, "增氧机开启 → 经过复核周期 → DO 达到恢复阈值 → 发送停机 → 事件关闭", 0.95, 5.2, 4.9, 0.32, size=11, color=INK, margin=0)
    pill(slide, "RESOLVED", 0.95, 5.56, 1.18, MINT, NAVY)
    text(slide, "处置结果、复核时间、设备回执均进入审计与日报", 2.3, 5.61, 3.5, 0.18, size=8.5, color=MUTED, margin=0)
    card(slide, 6.58, 4.52, 6.1, 1.42, PALE_CORAL, CORAL)
    text(slide, "自动动作失败怎么办？", 6.86, 4.78, 2.55, 0.26, size=15, color=CORAL, bold=True, margin=0)
    text(slide, "设备不健康 / MQTT 无回执 / 模型输出无效 / 超时 → 不硬写设备，转人工任务并保留失败原因。", 6.86, 5.18, 5.35, 0.35, size=10.5, color=INK, margin=0)
    pill(slide, "MANUAL_REQUIRED", 6.86, 5.56, 1.75, CORAL, WHITE)
    text(slide, "这条分支也是闭环的一部分", 8.85, 5.61, 2.8, 0.18, size=8.5, color=MUTED, margin=0)
    source_note(slide, "示例使用模拟传感器与本地 MQTT Broker；恢复阈值与复核周期可配置，演示不连接真实生产设备。")

    # 5. Actual action boundary
    slide = presentation.slides.add_slide(blank)
    base(slide, "Agent 如何真正做动作：设备写入有边界", "04 / 技术深度", number=5)
    text(slide, "大模型不直接改设备状态。它只产生结构化决策，确定性 Skill 再决定这条动作是否可执行。", 0.68, 1.78, 12, 0.4, size=14.5, color=INK, bold=True, margin=0)
    card(slide, 0.68, 2.5, 3.52, 3.65, PALE_BLUE, BLUE)
    text(slide, "大模型决策输出", 0.96, 2.78, 2.8, 0.26, size=16, color=BLUE, bold=True, margin=0)
    text(slide, "只展示业务必要字段，不暴露 API Key、供应商路由等配置。", 0.96, 3.15, 2.7, 0.38, size=9.5, color=MUTED, margin=0)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.96, 3.8, 2.96, 1.72, WHITE, rgb("C9DDF6"), True)
    text(slide, '{\n  "action": "EXECUTE",\n  "device_id": "aerator-b01-1",\n  "target_state": "on",\n  "risk": "L1",\n  "verification_delay": 300\n}', 1.16, 4.0, 2.55, 1.35, size=10.5, color=INK, font="Noto Sans Mono CJK SC", margin=0)
    line(slide, 4.3, 4.3, 4.73, 4.3, BLUE, 2.2, True)
    card(slide, 4.78, 2.5, 3.52, 3.65, PALE_YELLOW, YELLOW)
    text(slide, "DeviceControlSkill", 5.06, 2.78, 2.8, 0.26, size=16, color=INK, bold=True, margin=0)
    bullet_list(slide, ["动作白名单", "设备能力与健康状态", "风险等级与审批规则", "幂等键与超时", "失败原因分流人工"], 5.04, 3.36, 2.82, 1.48, size=10.5)
    pill(slide, "允许 → 发布", 5.06, 5.28, 1.18, MINT, NAVY)
    pill(slide, "拒绝 → 人工", 6.38, 5.28, 1.18, CORAL, WHITE)
    line(slide, 8.4, 4.3, 8.82, 4.3, YELLOW, 2.2, True)
    card(slide, 8.88, 2.5, 3.77, 3.65, PALE_MINT, MINT)
    text(slide, "MQTT 设备边界", 9.16, 2.78, 2.8, 0.26, size=16, color=INK, bold=True, margin=0)
    text(slide, "发布主题", 9.16, 3.32, 1.2, 0.2, size=10, color=MUTED, bold=True, margin=0)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.16, 3.58, 3.1, 0.52, WHITE, rgb("C5E8D8"), True)
    text(slide, "fishagent/ponds/B-01/devices/aerator-b01-1/commands", 9.32, 3.76, 2.78, 0.16, size=7.4, color=INK, font="Noto Sans Mono CJK SC", margin=0)
    text(slide, "回执与状态", 9.16, 4.4, 1.2, 0.2, size=10, color=MUTED, bold=True, margin=0)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.16, 4.68, 3.1, 0.82, WHITE, rgb("C5E8D8"), True)
    text(slide, "命令：开启增氧机\n回执：ACKNOWLEDGED\n状态：shadow_state = on", 9.32, 4.78, 2.78, 0.6, size=9, color=INK, margin=0)
    source_note(slide, "代码对应：fishagent.agent_runtime.skills.device_control.DeviceControlSkill；动作发布与模拟设备均通过 MQTT。")

    # 6. Multimodal
    slide = presentation.slides.add_slide(blank)
    base(slide, "多模态不是装饰：异常事实要被交叉验证", "05 / 多模态能力", number=6)
    text(slide, "当传感器读数不可信或异常原因不唯一时，Agent 将图像、天气、知识库和设备状态合并成可解释证据。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    add_picture(slide, ASSET / "b01-surface.png", 0.68, 2.42, 3.35, 2.1)
    add_picture(slide, ASSET / "b01-underwater.png", 4.2, 2.42, 3.35, 2.1)
    text(slide, "水面摄像头 · 浮头 / 水面状态", 0.82, 4.6, 3.0, 0.2, size=9, color=INK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    text(slide, "水下摄像头 · 鱼群活动 / 体表观察", 4.34, 4.6, 3.0, 0.2, size=9, color=INK, bold=True, align=PP_ALIGN.CENTER, margin=0)
    card(slide, 7.78, 2.42, 4.9, 2.4, WHITE, LINE)
    text(slide, "证据融合", 8.08, 2.7, 2.1, 0.26, size=16, color=INK, bold=True, margin=0)
    evidence = [
        ("传感器", "DO 2.8 mg/L / 质量 GOOD", BLUE),
        ("天气", "低气压 / 降雨前兆", CYAN),
        ("知识库", "低氧处置参考与风险提示", PURPLE),
        ("设备", "增氧机健康 / 在线状态", MINT),
    ]
    for idx, (label, body, accent) in enumerate(evidence):
        y = 3.16 + idx * 0.37
        pill(slide, label, 8.08, y, 0.72, accent, WHITE)
        text(slide, body, 8.98, y + 0.05, 3.25, 0.18, size=9.5, color=MUTED, margin=0)
    card(slide, 7.78, 5.08, 4.9, 1.05, PALE_CORAL, CORAL)
    text(slide, "安全分流", 8.08, 5.31, 1.2, 0.22, size=13, color=CORAL, bold=True, margin=0)
    text(slide, "传感器 SUSPECT / 设备不健康 / 高风险判断 → 不自动确诊，不自动投药，创建详细人工任务。", 9.28, 5.25, 3.05, 0.44, size=9.5, color=INK, margin=0)
    source_note(slide, "图片为演示环境的模拟摄像头样例；未使用真实养殖户个人信息或生产现场视频。")

    # 7. Product evidence
    slide = presentation.slides.add_slide(blank)
    base(slide, "产品体验：评委可以从页面验证每一步", "06 / 可演示产品", number=7)
    add_picture(slide, monitor_top, 0.68, 1.78, 8.02, 5.08)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 1.78, 8.02, 5.08, WHITE, LINE, False).fill.transparency = 1
    numbered_callout(slide, 1, "主动巡塘", "先向传感器发布 REPORT_NOW，再分析这批新鲜数据。", 9.05, 1.84, 3.55, PALE_BLUE, BLUE)
    numbered_callout(slide, 2, "水质监控面板", "异常来源高亮，所有 7 个指标可逐池塘查看。", 9.05, 2.72, 3.55, PALE_YELLOW, YELLOW)
    numbered_callout(slide, 3, "悬浮告警胶囊", "活跃告警固定在左侧，点击进入告警流程。", 9.05, 3.6, 3.55, PALE_CORAL, CORAL)
    numbered_callout(slide, 4, "Agent 执行中心", "展示待处理、进行中、已完成，而不是静态流程图。", 9.05, 4.48, 3.55, PALE_MINT, MINT)
    numbered_callout(slide, 5, "聊天与报告", "对话可查询快照，日报可下载独立 HTML。", 9.05, 5.36, 3.55, rgb("F0EDFF"), PURPLE)
    source_note(slide, "截图：本地运行实例 http://127.0.0.1:3000；此图为正常基线，故显示 28/28，故障 Demo 注入后会显示 27/28。")

    # 8. Verifiable evidence
    slide = presentation.slides.add_slide(blank)
    base(slide, "核心功能有证据：页面、接口、日志和报告彼此对得上", "07 / 可验证材料", number=8)
    text(slide, "评委可以从真实运行快照、报告预览、HTTP 接口和回归测试四个层面复核同一条任务链。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    add_picture(slide, report_top, 0.68, 2.38, 7.6, 4.75)
    numbered_callout(slide, 1, "真实报告正文", "趋势图来自当前快照，不是 PPT 绘制的示意图。", 8.62, 2.42, 3.98, PALE_BLUE, BLUE)
    numbered_callout(slide, 2, "自动日报与交付", "每日 23:59 自动生成；历史版本、下载 HTML、删除报告均有页面入口。", 8.62, 3.3, 3.98, PALE_MINT, MINT)
    numbered_callout(slide, 3, "接口可复现", "Demo、Agent steps、MQTT 命令、报告均有 HTTP API。", 8.62, 4.18, 3.98, PALE_YELLOW, YELLOW)
    numbered_callout(slide, 4, "测试可回归", "当前工程回归：106 passed，2 skipped；另有 ruff 检查。", 8.62, 5.06, 3.98, PALE_CORAL, CORAL)
    card(slide, 8.62, 6.05, 3.98, 0.68, NAVY, NAVY)
    text(slide, "复现入口：POST /api/v1/demo/success\n交付入口：GET /api/v1/reports/{id}/download", 8.86, 6.17, 3.5, 0.4, size=8.2, color=WHITE, font="Noto Sans Mono CJK SC", margin=0)
    source_note(slide, "截图：当前本地运行实例；报告包含真实快照趋势图、告警、自动操作、人工任务和设备操作日志。")

    # 9. Data flow / architecture
    slide = presentation.slides.add_slide(blank)
    base(slide, "模型、Agent 与工具接口：谁负责什么", "08 / 工程架构", number=9)
    text(slide, "系统把每一次输入、模型消息、决策、Skill 调用、设备回执和复核结果写入可观察轨迹。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    lanes = [
        ("感知层", "传感器 / 摄像头 / 天气", BLUE, PALE_BLUE),
        ("消息层", "本地 MQTT Broker\n上报 + 控制", CYAN, PALE_BLUE),
        ("Agent 层", "CrewAI Supervisor\n专职 Agent 协作", PURPLE, rgb("F0EDFF")),
        ("工具层", "RAG / Weather /\nDeviceControlSkill", YELLOW, PALE_YELLOW),
        ("执行层", "MQTT 控制消息\n模拟设备回执", MINT, PALE_MINT),
        ("验证层", "主动巡塘 / 复核\n告警状态更新", CORAL, PALE_CORAL),
        ("交付层", "页面 / 聊天 / 日报\n审计与人工任务", BLUE, PALE_BLUE),
    ]
    for idx, (label, body, accent, fill) in enumerate(lanes):
        left = 0.68 + idx * 1.78
        card(slide, left, 2.48, 1.5, 1.52, fill, accent)
        text(slide, label, left + 0.12, 2.75, 1.26, 0.24, size=11, color=INK, bold=True, align=PP_ALIGN.CENTER, margin=0)
        text(slide, body, left + 0.12, 3.16, 1.26, 0.48, size=8.2, color=MUTED, align=PP_ALIGN.CENTER, margin=0)
        if idx < len(lanes) - 1:
            line(slide, left + 1.52, 3.23, left + 1.72, 3.23, accent, 1.5, True)
    card(slide, 0.68, 4.32, 7.05, 2.02, WHITE, LINE)
    text(slide, "Agent 角色分工", 0.96, 4.57, 2.0, 0.24, size=15, color=INK, bold=True, margin=0)
    roles = [
        ("supervisor-agent", "校验触发目标，委派 Agent，汇总模型结果", PURPLE),
        ("sensor-monitor-agent", "请求 REPORT_NOW，读取池塘新鲜快照", BLUE),
        ("patrol-analysis-agent", "分析水质、设备影子、天气与活跃告警", CYAN),
        ("action-planning-agent", "检索知识库，形成 action / rationale", YELLOW),
        ("execution / verification", "Skill 下发 MQTT；复核、重试或转人工", MINT),
    ]
    for idx, (agent, responsibility, accent) in enumerate(roles):
        y = 4.92 + idx * 0.27
        add_shape(slide, MSO_SHAPE.OVAL, 0.98, y + 0.02, 0.16, 0.16, accent, accent, False)
        text(slide, agent, 1.25, y, 2.25, 0.18, size=8.7, color=INK, bold=True, font="Noto Sans Mono CJK SC", margin=0)
        text(slide, responsibility, 3.62, y, 3.65, 0.18, size=8.7, color=MUTED, margin=0)
    card(slide, 7.94, 4.32, 4.74, 2.02, PALE_YELLOW, YELLOW)
    text(slide, "工具接口与协议", 8.22, 4.57, 2.2, 0.24, size=15, color=INK, bold=True, margin=0)
    tools = [
        "只读  get_pond_snapshot(pond_id)",
        "只读  get_weather_context(pond_id)",
        "RAG   search_knowledge_base(query, species, metric)",
        "只读  get_device_shadow_state(pond_id)",
        "写入  DeviceControlSkill → MQTT set_state",
    ]
    for idx, item in enumerate(tools):
        text(slide, item, 8.22, 4.94 + idx * 0.26, 4.05, 0.18, size=8.2, color=INK if idx == 4 else MUTED, bold=idx == 4, font="Noto Sans Mono CJK SC", margin=0)
    card(slide, 0.68, 6.58, 12.0, 0.38, NAVY, NAVY)
    text(slide, "协议：传感器 commands → REPORT_NOW；设备 fishagent/ponds/{pond_id}/devices/{device_id}/commands → set_state；证据落地 agent_runs / events / commands / verification_results", 0.9, 6.68, 11.55, 0.18, size=8.1, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)

    # 10. Safety
    slide = presentation.slides.add_slide(blank)
    base(slide, "安全边界：让 Agent 能做事，也知道不能做什么", "09 / 安全与合规", number=10)
    text(slide, "水产养殖涉及生物、设备和现场操作风险。主要风险是误判、隐私泄露和误动作；系统用策略门、人工接管、最小化数据和审计来控制。", 0.68, 1.78, 12, 0.42, size=13.2, color=INK, bold=True, margin=0)
    safety = [
        ("模型层", "负责理解\n规划与解释", "不直接写设备", PALE_BLUE, BLUE),
        ("策略门", "白名单 / 能力\n风险 / 幂等 / 审批", "Skill 才能执行", PALE_YELLOW, YELLOW),
        ("人工路由", "无效输出 / 超时\n回执失败 / 高风险", "详细任务接管", PALE_CORAL, CORAL),
        ("审计层", "证据 / 决策 / 回执\n复核 / 报告", "每一步可追溯", PALE_MINT, MINT),
    ]
    for idx, (title, body, note, fill, accent) in enumerate(safety):
        left = 0.68 + idx * 3.03
        card(slide, left, 2.5, 2.7, 1.75, fill, accent)
        text(slide, title, left + 0.22, 2.79, 2.25, 0.26, size=16, color=INK, bold=True, margin=0)
        text(slide, body, left + 0.22, 3.2, 2.2, 0.48, size=11, color=MUTED, margin=0)
        pill(slide, note, left + 0.22, 3.83, 1.75, accent, WHITE if accent in [BLUE, CORAL, MINT] else NAVY)
    card(slide, 0.68, 4.68, 5.82, 1.48, WHITE, LINE)
    text(slide, "当前演示数据授权与隐私状态", 0.96, 4.96, 3.1, 0.25, size=15, color=INK, bold=True, margin=0)
    bullet_list(slide, [
        "传感器、设备、天气、知识库是项目生成的模拟数据",
        "摄像头使用项目内模拟图片，不含真实人员或养殖户信息",
        "API Key 脱敏展示；公网部署必须开启认证并限制访问来源",
    ], 0.96, 5.36, 5.15, 0.65, size=9.5, gap=1)
    card(slide, 6.78, 4.68, 5.9, 1.48, PALE_CORAL, CORAL)
    text(slide, "真实接入前必须完成", 7.06, 4.96, 2.3, 0.25, size=15, color=CORAL, bold=True, margin=0)
    text(slide, "养殖场 / 养殖户书面授权、用途最小化、留存与删除策略\n第三方模型 / API 清单、数据传输范围、密钥管理\n现场负责人和专业人员确认，Agent 不替代最终判断", 7.06, 5.3, 5.1, 0.62, size=9.3, color=INK, margin=0)
    source_note(slide, "按手册 9.1—9.4：明确数据来源、授权、第三方依赖、风险提示和行业责任边界。当前 PPT 不把模拟数据包装成真实试点结果。")

    # 11. Reproducibility / open source
    slide = presentation.slides.add_slide(blank)
    base(slide, "可运行、可复现、可复用", "10 / 开放贡献", number=11)
    text(slide, "评委拿到的不是一张概念图，而是一套可以从零启动、注入事件、观察数据流并复盘结果的工程。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    card(slide, 0.68, 2.46, 5.45, 3.36, NAVY, NAVY)
    text(slide, "从零启动", 0.98, 2.76, 1.6, 0.28, size=17, color=WHITE, bold=True, margin=0)
    commands = [
        "uv sync --extra agent",
        "docker compose up -d --build",
        "http://<host>:3000",
        "POST /api/v1/demo/success",
        "uv run pytest -q",
    ]
    for idx, command in enumerate(commands):
        y = 3.28 + idx * 0.42
        pill(slide, f"0{idx + 1}", 0.98, y, 0.42, CYAN, NAVY)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.58, y, 3.9, 0.31, rgb("163D70"), rgb("2B6195"), True)
        text(slide, command, 1.76, y + 0.06, 3.55, 0.16, size=9.5, color=WHITE, font="Noto Sans Mono CJK SC", margin=0)
    text(slide, "当前回归：106 passed, 2 skipped", 0.98, 5.55, 4.4, 0.2, size=9.5, color=YELLOW, bold=True, margin=0)
    card(slide, 6.45, 2.46, 6.23, 3.36, WHITE, LINE)
    text(slide, "可以复用的开放组件", 6.74, 2.76, 3.2, 0.28, size=17, color=INK, bold=True, margin=0)
    components = [
        ("Agent 运行时", "CrewAI Supervisor + 专职 Agent", BLUE),
        ("设备 Skill", "确定性 DeviceControlSkill + MQTT", MINT),
        ("数据规格", "模拟传感器 / 摄像头 / 天气 / 知识库", PURPLE),
        ("交付模板", "日报 HTML、审计、人工任务、运行轨迹", CORAL),
        ("部署材料", "uv / Docker Compose / README / 测试", YELLOW),
    ]
    for idx, (label, body, accent) in enumerate(components):
        y = 3.28 + idx * 0.46
        add_shape(slide, MSO_SHAPE.OVAL, 6.76, y + 0.02, 0.25, 0.25, accent, accent, False)
        text(slide, label, 7.18, y, 1.35, 0.2, size=9.8, color=INK, bold=True, margin=0)
        text(slide, body, 8.72, y, 3.45, 0.2, size=9.3, color=MUTED, margin=0)
    text(slide, "GitHub：github.com/simonliu009/fishagent", 6.76, 5.54, 5.0, 0.2, size=9.5, color=BLUE, bold=True, margin=0)
    source_note(slide, "复赛 / 决赛需按组委会要求补充最终仓库、部署说明、许可证和第三方依赖清单。")

    # 12. Demo runbook
    slide = presentation.slides.add_slide(blank)
    base(slide, "给评委的 3 分钟现场 Demo 路径", "11 / 现场演示", number=12)
    text(slide, "从正常态开始，注入一个可重复的低溶氧事件；每一步都能在页面、Agent 轨迹或审计中找到证据。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    demo = [
        ("00:00", "正常态", "4 个池塘均有最新读数\n设备在线比例可见", BLUE, PALE_BLUE),
        ("00:20", "注入低溶氧", "点击 Demo：B-01 DO 降至 2.8\n页面中间出现注入反馈", CORAL, PALE_CORAL),
        ("00:45", "触发逐塘巡检", "正在触发逐塘巡检…\n主动请求 MQTT REPORT_NOW", CYAN, PALE_BLUE),
        ("01:20", "Agent 决策", "流式显示输入、回复、action\nEXECUTE → Skill 校验", PURPLE, rgb("F0EDFF")),
        ("01:50", "设备与告警", "发布 MQTT 控制\n告警流程逐步显示执行结果", MINT, PALE_MINT),
        ("02:30", "复核与交付", "下一次巡塘更新复核\n生成日报 / 聊天查询状态", YELLOW, PALE_YELLOW),
    ]
    for idx, (timecode, title, body, accent, fill) in enumerate(demo):
        left = 0.68 + (idx % 3) * 4.17
        top = 2.5 + (idx // 3) * 1.55
        card(slide, left, top, 3.72, 1.2, fill, accent)
        pill(slide, timecode, left + 0.22, top + 0.2, 0.72, accent, WHITE if accent not in [YELLOW] else NAVY)
        text(slide, title, left + 1.13, top + 0.19, 2.2, 0.23, size=13, color=INK, bold=True, margin=0)
        text(slide, body, left + 0.22, top + 0.62, 3.2, 0.37, size=9.5, color=MUTED, margin=0)
    card(slide, 0.68, 5.88, 7.7, 0.68, WHITE, LINE)
    text(slide, "现场入口", 0.94, 6.08, 0.9, 0.2, size=11, color=INK, bold=True, margin=0)
    text(slide, "http://<评审可访问主机>:3000   ·   账号 / 密码 / 录屏地址：待补", 1.95, 6.08, 6.0, 0.2, size=10, color=BLUE, font="Noto Sans Mono CJK SC", margin=0)
    card(slide, 8.65, 5.88, 4.03, 0.68, PALE_CORAL, CORAL)
    text(slide, "失败分支也要演示", 8.92, 6.08, 1.55, 0.2, size=10.5, color=CORAL, bold=True, margin=0)
    text(slide, "设备失败 → 人工任务 + 原因", 10.55, 6.08, 1.8, 0.2, size=9.5, color=INK, margin=0)
    source_note(slide, "演示入口与账号属于最终交付信息，当前保留占位符，避免把本地环境地址误当成比赛访问地址。")

    # 13. Evidence matrix
    slide = presentation.slides.add_slide(blank)
    base(slide, "用证据回答评委的评分表", "12 / 评审对照", number=13)
    text(slide, "每个评分维度对应一个可点击、可运行或可复盘的证据，不用技术名词堆叠替代结果。", 0.68, 1.78, 12, 0.4, size=14, color=INK, bold=True, margin=0)
    cols = [0.68, 2.48, 6.0, 9.22]
    widths = [1.55, 3.28, 2.95, 3.46]
    headers = ["评审维度", "我们展示什么", "当前工程证据", "封版前补齐"]
    for left, width, label in zip(cols, widths, headers):
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, 2.45, width, 0.42, NAVY, NAVY, False)
        text(slide, label, left + 0.08, 2.57, width - 0.16, 0.16, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
    rows = [
        ("行业场景价值 25%", "低溶氧处置、设备离线、人工接管", "四池塘 / 28 设备 / 7 指标；health Demo 可复现 27/28", "□ 真实试点收益数据"),
        ("Agent 闭环 25%", "理解 → 规划 → 工具 → 执行 → 复核", "CrewAI + Skill + MQTT + verification", "□ 现场操作确认"),
        ("产品体验 20%", "告警、聊天、轨迹、报告", "页面截图 + API + 106 条回归测试", "□ Demo 录屏链接"),
        ("技术深度 15%", "多模态、RAG、调度、持久化", "Docker / Postgres / Redis / MinIO", "□ 性能与模型成本"),
        ("安全合规 10%", "风险分级、策略门、人工任务", "模拟数据 / 审计 / 失败分流", "□ 授权与依赖清单"),
        ("开放复用 5%", "组件、规格、文档、部署", "GitHub + README + 测试", "□ License / 贡献指南"),
    ]
    for row_idx, row in enumerate(rows):
        top = 2.87 + row_idx * 0.61
        fill = WHITE if row_idx % 2 == 0 else rgb("EEF3F9")
        for col_idx, (left, width, value) in enumerate(zip(cols, widths, row)):
            add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, 0.61, fill, LINE, False)
            color = CORAL if col_idx == 3 else INK if col_idx == 0 else MUTED
            text(slide, value, left + 0.1, top + 0.11, width - 0.2, 0.38, size=8.7 if col_idx else 9.2, color=color, bold=col_idx == 0, margin=0)
    card(slide, 0.68, 6.65, 12.0, 0.34, PALE_YELLOW, YELLOW)
    text(slide, "红线意识：PPT 只能说明方案，Demo、日志、视频和可复现工程才是最终验证。", 0.92, 6.74, 11.5, 0.16, size=9.5, color=INK, bold=True, align=PP_ALIGN.CENTER, margin=0)

    # 14. Closing / placeholders
    slide = presentation.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    dot_field(slide, 8.65, 0.45, 3.8, 3.8)
    add_picture(slide, ASSET / "goai-logo.png", 10.8, 0.42, 1.32, 0.45)
    text(slide, "封版前，只补齐四件事", 0.72, 0.72, 6.5, 0.35, size=11, color=YELLOW, bold=True, margin=0)
    text(slide, "让一个可运行 Demo\n变成评委记得住的作品", 0.68, 1.32, 7.1, 1.12, size=30, color=WHITE, bold=True, margin=0)
    closing = [
        ("01", "团队与作者", "团队名称、成员、分工、联系方式"),
        ("02", "真实价值", "试点对象、收益指标、失败案例与边界"),
        ("03", "现场证据", "Demo 视频、访问地址、演示账号、运行记录"),
        ("04", "合规与开放", "授权、License、第三方依赖与数据说明"),
    ]
    for idx, (number, title, body) in enumerate(closing):
        left = 0.72 + (idx % 2) * 4.35
        top = 3.0 + (idx // 2) * 1.25
        card(slide, left, top, 3.85, 0.95, rgb("143A6B"), rgb("2C6296"))
        pill(slide, number, left + 0.22, top + 0.28, 0.48, CYAN, NAVY)
        text(slide, title, left + 0.9, top + 0.2, 2.55, 0.22, size=13, color=WHITE, bold=True, margin=0)
        text(slide, body, left + 0.9, top + 0.5, 2.65, 0.2, size=8.8, color=rgb("B8CAE0"), margin=0)
    text(slide, "智渔Agent 2.0 = 养殖场的感知、决策、执行、复核协同层", 0.72, 6.12, 8.8, 0.3, size=15, color=CYAN, bold=True, margin=0)
    text(slide, "谢谢 / Q&A", 0.72, 6.62, 3.0, 0.35, size=18, color=WHITE, bold=True, margin=0)
    text(slide, "参赛材料依据：GOAI 无界应用｜Boundless Agents 参赛手册；官网：goaihz.com", 0.72, 7.15, 8.4, 0.18, size=7.5, color=rgb("B8CAE0"), margin=0)

    presentation.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_slides()
